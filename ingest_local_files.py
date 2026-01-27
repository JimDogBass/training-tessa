"""
Training Tessa - Local File Ingestion Script (Enhanced)
Reads training documents, extracts text and images, creates embeddings,
and stores directly in Supabase with source document names.
Now using Gemini for embeddings and image descriptions.
"""

import os
import sys
import json
import time
import base64
import httpx
from pathlib import Path
from io import BytesIO
from google import genai

# Configuration
TRAINING_FOLDER = r"C:\Users\JoelBentley\OneDrive - Meraki Talent\Agent Content\Training"

# Gemini API (for embeddings and image descriptions)
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY")
GEMINI_EMBEDDING_MODEL = "text-embedding-004"
GEMINI_VISION_MODEL = "gemini-2.5-flash"

# Retry settings
MAX_RETRIES = 3
RETRY_DELAY = 5  # seconds

if not GEMINI_API_KEY:
    print("ERROR: GEMINI_API_KEY environment variable not set.")
    print("Set it with: set GEMINI_API_KEY=your-api-key")
    sys.exit(1)

# Configure Gemini client
gemini_client = genai.Client(api_key=GEMINI_API_KEY)

# Supabase
SUPABASE_URL = os.environ.get("SUPABASE_URL", "https://uwxsflcpaigcygfhxzzl.supabase.co")
SUPABASE_KEY = os.environ.get("SUPABASE_KEY")

if not SUPABASE_KEY:
    print("ERROR: SUPABASE_KEY environment variable not set.")
    print("Set it with: set SUPABASE_KEY=your-supabase-key")
    sys.exit(1)

# Chunking settings
CHUNK_SIZE = 1000  # characters per chunk
CHUNK_OVERLAP = 200  # overlap between chunks

REQUEST_TIMEOUT = 120


def extract_docx(file_path):
    """Extract text and images from .docx files."""
    try:
        from docx import Document
        from docx.opc.constants import RELATIONSHIP_TYPE as RT

        doc = Document(file_path)
        text_parts = []
        images = []

        # Extract text
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # Extract images
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_data = rel.target_part.blob
                    images.append(image_data)
                except:
                    pass

        return '\n\n'.join(text_parts), images
    except Exception as e:
        print(f"  Error extracting docx: {e}")
        return None, []


def extract_pptx(file_path):
    """Extract text and images from .pptx files."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(file_path)
        text_parts = []
        images = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                # Extract text
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

                # Extract images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        images.append(image.blob)
                    except:
                        pass

            if slide_text:
                text_parts.append(f"Slide {slide_num}:\n" + '\n'.join(slide_text))

        return '\n\n'.join(text_parts), images
    except Exception as e:
        print(f"  Error extracting pptx: {e}")
        return None, []


def extract_pdf(file_path):
    """Extract text and images from .pdf files."""
    try:
        import pdfplumber
        from PIL import Image

        text_parts = []
        images = []

        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                # Extract text
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    text_parts.append(f"Page {page_num}:\n{page_text}")

                # Extract images
                try:
                    for img in page.images:
                        # pdfplumber gives image info but not always the raw data
                        pass
                except:
                    pass

        return '\n\n'.join(text_parts), images
    except Exception as e:
        print(f"  Error extracting pdf: {e}")
        return None, []


def extract_odt(file_path):
    """Extract text from .odt files."""
    try:
        from odf import text as odf_text
        from odf.opendocument import load
        doc = load(file_path)
        all_text = []
        for elem in doc.getElementsByType(odf_text.P):
            text_content = ''.join(node.data for node in elem.childNodes if hasattr(node, 'data'))
            if text_content.strip():
                all_text.append(text_content)
        return '\n\n'.join(all_text), []
    except Exception as e:
        print(f"  Error extracting odt: {e}")
        return None, []


def extract_content(file_path):
    """Extract text and images based on file extension."""
    ext = Path(file_path).suffix.lower()

    if ext in ['.docx', '.doc']:
        return extract_docx(file_path)
    elif ext in ['.pptx', '.ppt']:
        return extract_pptx(file_path)
    elif ext == '.pdf':
        return extract_pdf(file_path)
    elif ext in ['.odt', '.odp']:
        return extract_odt(file_path)
    elif ext == '.txt':
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read(), []
    else:
        print(f"  Unsupported file type: {ext}")
        return None, []


def describe_image_with_vision(image_data, filename):
    """Use Gemini Vision to describe an image."""
    try:
        from PIL import Image
        from io import BytesIO

        # Convert bytes to PIL Image
        image = Image.open(BytesIO(image_data))

        prompt = f"""You are a helpful assistant that describes images from training documents.
Please describe this image from the training document '{filename}'.
Include any visible text, diagrams, charts, or key information.
Focus on information that would be useful for training purposes.
Keep the description concise but complete."""

        response = gemini_client.models.generate_content(
            model=GEMINI_VISION_MODEL,
            contents=[prompt, image]
        )
        description = response.text
        return f"[Image Description]: {description}"
    except Exception as e:
        print(f"    Error describing image: {e}")
        return None


def chunk_text(text, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    """Split text into overlapping chunks."""
    if len(text) <= chunk_size:
        return [text]

    chunks = []
    start = 0

    while start < len(text):
        end = start + chunk_size

        # Try to break at a sentence or paragraph
        if end < len(text):
            # Look for paragraph break
            newline_pos = text.rfind('\n\n', start, end)
            if newline_pos > start + chunk_size // 2:
                end = newline_pos
            else:
                # Look for sentence break
                for sep in ['. ', '! ', '? ', '\n']:
                    sep_pos = text.rfind(sep, start, end)
                    if sep_pos > start + chunk_size // 2:
                        end = sep_pos + len(sep)
                        break

        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)

        start = end - overlap
        if start < 0:
            start = 0

    return chunks


def create_embedding(text):
    """Create embedding using Gemini with retry logic."""
    for attempt in range(MAX_RETRIES):
        try:
            result = gemini_client.models.embed_content(
                model=GEMINI_EMBEDDING_MODEL,
                contents=text,
                config={"task_type": "RETRIEVAL_DOCUMENT"}
            )
            return result.embeddings[0].values
        except Exception as e:
            if attempt < MAX_RETRIES - 1:
                print(f"    Embedding attempt {attempt + 1} failed, retrying in {RETRY_DELAY}s...")
                time.sleep(RETRY_DELAY)
            else:
                print(f"    Error creating embedding after {MAX_RETRIES} attempts: {e}")
                return None


def insert_to_supabase(content, embedding, source_document):
    """Insert document chunk into Supabase."""
    try:
        url = f"{SUPABASE_URL}/rest/v1/documents"

        with httpx.Client(timeout=REQUEST_TIMEOUT) as client:
            response = client.post(
                url,
                headers={
                    "apikey": SUPABASE_KEY,
                    "Authorization": f"Bearer {SUPABASE_KEY}",
                    "Content-Type": "application/json",
                    "Prefer": "return=minimal"
                },
                json={
                    "content": content,
                    "embedding": embedding,
                    "source_document": source_document
                }
            )
            response.raise_for_status()
            return True
    except Exception as e:
        print(f"    Error inserting to Supabase: {e}")
        return False


def clear_supabase_documents():
    """Clear all existing documents from Supabase."""
    try:
        url = f"{SUPABASE_URL}/rest/v1/documents"

        with httpx.Client(timeout=REQUEST_TIMEOUT) as client:
            # Delete all documents
            response = client.delete(
                url,
                headers={
                    "apikey": SUPABASE_KEY,
                    "Authorization": f"Bearer {SUPABASE_KEY}",
                    "Content-Type": "application/json"
                },
                params={"id": "gt.0"}  # Delete all where id > 0
            )
            response.raise_for_status()
            print("Cleared all existing documents from Supabase")
            return True
    except Exception as e:
        print(f"Error clearing Supabase: {e}")
        return False


def process_file(file_path, filename, process_images=True):
    """Process a single file: extract content, chunk, embed, and store."""
    print(f"Processing: {filename}")

    # Extract text and images
    text, images = extract_content(file_path)

    if not text or len(text.strip()) < 50:
        print(f"  Skipped - no text extracted or too short")
        return 0, 0

    print(f"  Extracted {len(text)} characters, {len(images)} images")

    # Process images with Gemini Vision
    image_descriptions = []
    if process_images and images:
        print(f"  Processing {len(images)} images with Gemini Vision...")
        for i, img_data in enumerate(images[:5]):  # Limit to 5 images per doc
            desc = describe_image_with_vision(img_data, filename)
            if desc:
                image_descriptions.append(desc)
                print(f"    Image {i+1} described")
            time.sleep(1)  # Rate limiting

    # Combine text with image descriptions
    full_content = text
    if image_descriptions:
        full_content += "\n\n" + "\n\n".join(image_descriptions)

    # Chunk the content
    chunks = chunk_text(full_content)
    print(f"  Split into {len(chunks)} chunks")

    # Process each chunk
    success_count = 0
    for i, chunk in enumerate(chunks):
        # Create embedding
        embedding = create_embedding(chunk)
        if not embedding:
            continue

        # Insert to Supabase
        if insert_to_supabase(chunk, embedding, filename):
            success_count += 1

        # Small delay for rate limiting
        time.sleep(0.5)

    print(f"  [OK] Stored {success_count}/{len(chunks)} chunks")
    return success_count, len(images)


def main():
    """Main ingestion function."""
    print("=" * 60)
    print("Training Tessa - Enhanced Local File Ingestion")
    print("=" * 60)
    print(f"\nSource folder: {TRAINING_FOLDER}")
    print(f"Supabase URL: {SUPABASE_URL}")

    # Check folder exists
    if not os.path.exists(TRAINING_FOLDER):
        print(f"ERROR: Folder not found: {TRAINING_FOLDER}")
        sys.exit(1)

    # Ask about clearing existing documents
    print("\n" + "-" * 60)
    print("Do you want to clear existing documents from Supabase first?")
    print("This will remove all current documents and re-ingest fresh.")
    print("-" * 60)

    response = input("Clear existing documents? (y/n): ").strip().lower()
    if response == 'y':
        clear_supabase_documents()

    # Ask about processing images
    print("\n" + "-" * 60)
    response = input("Process images with Gemini Vision? (y/n): ").strip().lower()
    process_images = response == 'y'
    print("-" * 60 + "\n")

    # Get all files
    files = [f for f in os.listdir(TRAINING_FOLDER) if os.path.isfile(os.path.join(TRAINING_FOLDER, f))]

    # Filter supported file types
    supported_extensions = ['.docx', '.doc', '.pptx', '.ppt', '.pdf', '.odt', '.odp', '.txt']
    files = [f for f in files if Path(f).suffix.lower() in supported_extensions]

    print(f"Found {len(files)} supported files to process\n")

    # Process each file
    total_chunks = 0
    total_images = 0
    success_files = 0

    for i, filename in enumerate(files, 1):
        file_path = os.path.join(TRAINING_FOLDER, filename)
        print(f"\n[{i}/{len(files)}] ", end="")

        chunks, images = process_file(file_path, filename, process_images)

        if chunks > 0:
            success_files += 1
            total_chunks += chunks
            total_images += images

    # Summary
    print("\n" + "=" * 60)
    print("INGESTION COMPLETE")
    print("=" * 60)
    print(f"Files processed: {success_files}/{len(files)}")
    print(f"Total chunks stored: {total_chunks}")
    print(f"Total images processed: {total_images}")
    print("\nTraining Tessa is ready to answer questions!")


if __name__ == "__main__":
    main()
